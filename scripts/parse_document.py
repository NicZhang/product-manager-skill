#!/usr/bin/env python3
"""Unified document parser for .docx, .xlsx, and .pptx files. Outputs structured markdown to stdout."""

import sys
from pathlib import Path


def parse_docx(file_path: str) -> str:
    from docx import Document
    from docx.oxml.ns import qn

    doc = Document(file_path)
    lines = []

    heading_map = {
        "Heading 1": "#",
        "Heading 2": "##",
        "Heading 3": "###",
    }

    for block in doc.element.body:
        tag = block.tag.split("}")[-1] if "}" in block.tag else block.tag

        if tag == "p":
            from docx.text.paragraph import Paragraph
            para = Paragraph(block, doc)
            text = para.text.strip()
            if not text:
                continue
            style_name = para.style.name if para.style else ""
            if style_name in heading_map:
                lines.append(f"{heading_map[style_name]} {text}")
            elif style_name in ("List Paragraph", "List Bullet", "List Number"):
                lines.append(f"- {text}")
            else:
                lines.append(text)

        elif tag == "tbl":
            from docx.table import Table
            table = Table(block, doc)
            rows = table.rows
            if not rows:
                continue
            header_cells = [cell.text.strip().replace("|", "\\|") for cell in rows[0].cells]
            lines.append("| " + " | ".join(header_cells) + " |")
            lines.append("| " + " | ".join(["---"] * len(header_cells)) + " |")
            for row in rows[1:]:
                cells = [cell.text.strip().replace("|", "\\|") for cell in row.cells]
                lines.append("| " + " | ".join(cells) + " |")

    return "\n".join(lines)


def parse_xlsx(file_path: str) -> str:
    import openpyxl

    MAX_COLS = 20
    MAX_ROWS = 50

    wb = openpyxl.load_workbook(file_path, data_only=True)
    sections = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        lines = [f"## {sheet_name}"]

        all_rows = list(ws.iter_rows(values_only=True))
        non_empty_rows = [r for r in all_rows if any(c is not None for c in r)]
        if not non_empty_rows:
            continue

        total_rows = len(non_empty_rows)
        total_cols = max(len(r) for r in non_empty_rows)

        display_rows = non_empty_rows[:MAX_ROWS]
        display_cols = min(total_cols, MAX_COLS)

        def fmt(val):
            if val is None:
                return ""
            return str(val).strip().replace("|", "\\|")

        header = display_rows[0]
        header_cells = [fmt(header[i]) if i < len(header) else "" for i in range(display_cols)]
        lines.append("| " + " | ".join(header_cells) + " |")
        lines.append("| " + " | ".join(["---"] * display_cols) + " |")

        for row in display_rows[1:]:
            cells = [fmt(row[i]) if i < len(row) else "" for i in range(display_cols)]
            lines.append("| " + " | ".join(cells) + " |")

        summary_parts = [f"统计：{total_rows}行 × {total_cols}列"]
        if total_rows > MAX_ROWS:
            summary_parts.append(f"（仅显示前{MAX_ROWS}行）")
        if total_cols > MAX_COLS:
            summary_parts.append(f"（仅显示前{MAX_COLS}列）")
        lines.append("")
        lines.append("".join(summary_parts))

        sections.append("\n".join(lines))

    return "\n\n".join(sections)


def parse_pptx(file_path: str) -> str:
    from pptx import Presentation
    from pptx.util import Pt
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(file_path)
    sections = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        lines = [f"## Slide {slide_num}"]

        title_shape = slide.shapes.title
        if title_shape and title_shape.has_text_frame:
            title_text = title_shape.text.strip()
            if title_text:
                lines.append(f"**{title_text}**")

        for shape in slide.shapes:
            if shape == title_shape:
                continue

            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if not text:
                        continue
                    level = para.level or 0
                    indent = "  " * level
                    lines.append(f"{indent}- {text}")

            elif shape.has_table:
                table = shape.table
                rows = table.rows
                if not rows:
                    continue
                header_cells = [cell.text.strip().replace("|", "\\|") for cell in rows[0].cells]
                lines.append("| " + " | ".join(header_cells) + " |")
                lines.append("| " + " | ".join(["---"] * len(header_cells)) + " |")
                for row in rows[1:]:
                    cells = [cell.text.strip().replace("|", "\\|") for cell in row.cells]
                    lines.append("| " + " | ".join(cells) + " |")

        sections.append("\n".join(lines))

    return "\n\n".join(sections)


def parse_file(file_path: str) -> str:
    path = Path(file_path)

    if not path.exists():
        return f"错误：文件不存在 - {file_path}"

    ext = path.suffix.lower()
    supported = {".docx", ".doc", ".xlsx", ".xls", ".pptx", ".ppt"}

    if ext not in supported:
        return f"错误：不支持的文件格式 - {ext}\n支持格式：{', '.join(sorted(supported))}"

    try:
        if ext in (".docx", ".doc"):
            return parse_docx(file_path)
        elif ext in (".xlsx", ".xls"):
            return parse_xlsx(file_path)
        elif ext in (".pptx", ".ppt"):
            return parse_pptx(file_path)
    except Exception as e:
        return f"错误：解析文件失败 - {e}"


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("用法：python parse_document.py <file_path>", file=sys.stderr)
        sys.exit(1)

    result = parse_file(sys.argv[1])
    print(result)
